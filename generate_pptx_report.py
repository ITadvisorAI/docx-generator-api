from pptx import Presentation
from pptx.util import Inches
import os

def generate_pptx_report_with_charts(session_id):
    """
    Generates an executive PowerPoint report by embedding infrastructure charts
    into slides. If chart files are missing, warning text is added instead.
    """
    # Paths
    session_path = f"temp_sessions/{session_id}"
    chart_path = os.path.join(session_path, "charts")
    output_path = os.path.join(session_path, f"IT_Infrastructure_Executive_Report_{session_id}.pptx")
    template_path = "templates/Executive_PPT_Template.pptx"

    # Ensure session directory exists
    os.makedirs(session_path, exist_ok=True)

    # Load template or start from scratch
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Define charts to embed
    chart_files = [
        ("HW Tier Distribution", "hw_tier_distribution.png"),
        ("HW Environment Distribution", "hw_environment_distribution.png"),
        ("HW Device Type vs Tier", "hw_device_type_vs_tier.png"),
        ("SW Tier Distribution", "sw_tier_distribution.png"),
        ("SW Environment Distribution", "sw_environment_distribution.png")
    ]

    # Add a new slide per chart
    for title, filename in chart_files:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        if slide.shapes.title:
            slide.shapes.title.text = title

        chart_file = os.path.join(chart_path, filename)
        if os.path.exists(chart_file):
            slide.shapes.add_picture(chart_file, Inches(1), Inches(1.5), width=Inches(7))
        else:
            left = Inches(1)
            top = Inches(2)
            width = Inches(6)
            height = Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = f"⚠️ {title} chart not found."

    # Save final presentation
    prs.save(output_path)
    return output_path
