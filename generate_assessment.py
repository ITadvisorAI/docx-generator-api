import os
import re
import requests
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from drive_utils import upload_to_drive

# Section titles for Table of Contents (20 sections as per template)
SECTION_TITLES = [
    "Executive Summary",
    "Organization IT Landscape Overview",
    "Inventory Breakdown – Hardware",
    "Inventory Breakdown – Software",
    "Classification Tier Distribution",
    "Hardware Lifecycle Status",
    "Software Licensing and Compliance",
    "Security Posture and Vulnerabilities",
    "Performance Bottlenecks & Uptime Metrics",
    "System Reliability & Failover Readiness",
    "Scalability & Elasticity Opportunities",
    "Legacy Systems and Technical Debt",
    "Obsolete and High-Risk Platforms",
    "Cloud Migration Potential (Workload Mapping)",
    "Strategic Alignment of IT Assets",
    "Business Impact Analysis of Current Gaps",
    "Financial Implications – Cost of Obsolescence",
    "Environmental Impact and Sustainability",
    "Recommendations for Remediation & Upgrade",
    "Proposed Next Steps and Roadmap"
]

def build_table_of_contents(data: dict) -> str:
    """
    Generate a complete Table of Contents for all sections.
    """
    lines = []
    for idx, title in enumerate(SECTION_TITLES, start=1):
        lines.append(f"{idx}. {title}")
    # ✅ Fix: join with newline so each entry is on its own line :contentReference[oaicite:0]{index=0}
    return "\n".join(lines)

# Base paths
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_DOCX = os.path.join(BASE_DIR, "templates", "IT_Current_Status_Assessment_Report_Template.docx")
TEMPLATE_PPTX = os.path.join(BASE_DIR, "templates", "IT_Current_Status_Executive_Report_Template.pptx")
OUTPUT_ROOT = os.path.join(BASE_DIR, "temp_sessions")


def _to_direct_drive_url(url: str) -> str:
    # ... (unchanged) :contentReference[oaicite:1]{index=1}
    match = re.search(r"[?&]id=([\w\-]+)", url)
    if match:
        return f"https://drive.google.com/uc?export=download&id={match.group(1)}"
    match = re.search(r"/d/([\w\-]+)", url)
    if match:
        return f"https://drive.google.com/uc?export=download&id={match.group(1)}"
    return url


def generate_assessment_docs(**data):
    # Align keys & build TOC
    data["hw_gap_url"] = data.get("file_1_drive_url", "")
    data["sw_gap_url"] = data.get("file_2_drive_url", "")
    if "chart_paths" not in data:
        data["chart_paths"] = {k: v for k, v in data.items() if k.endswith("_chart")}
    data["table_of_contents"] = build_table_of_contents(data)

    # Session metadata
    session_id = data.get("session_id", "")
    report_date = data.get("report_date", "")
    print(f"[DEBUG] Generating docs for session: {session_id}", flush=True)

    # Section titles placeholders
    for idx, title in enumerate(SECTION_TITLES, start=1):
        data[f"section_{idx}_title"] = title

    # Map narratives into slide placeholders
    section_to_slide_map = {
        'executive_summary': 1,
        'it_landscape_overview': 2,
        'hardware_analysis': 3,
        'software_analysis': 4,
        'tier_classification_summary': 5,
        'hardware_lifecycle_chart': 6,
        'software_licensing_review': 7,
        'security_vulnerability_heatmap': 8,
        'performance_&_uptime_trends': 9,
        'system_reliability_overview': 10,
        'scalability_insights': 11,
        'legacy_system_exposure': 12,
        'obsolete_platform_matrix': 13,
        'cloud_migration_targets': 14,
        'strategic_it_alignment': 15,
        'business_impact_of_gaps': 16,
        'cost_of_obsolescence': 17,
        'environmental_impact_and_sustainability': 18,     # matches SECTION_TITLES index 18 :contentReference[oaicite:2]{index=2}
        'remediation_recommendations': 19,
        'roadmap_&_next_steps': 20
    }
    for key, sec_num in section_to_slide_map.items():
        data[f"slide_{key}"] = data.get(f"content_{sec_num}", "")

    # Prepare output dir & download charts
    session_dir = os.path.join(OUTPUT_ROOT, session_id)
    os.makedirs(session_dir, exist_ok=True)
    local_charts = {}
    for name, url in data.get("chart_paths", {}).items():
        try:
            dl = _to_direct_drive_url(url)
            r = requests.get(dl); r.raise_for_status()
            path = os.path.join(session_dir, f"{name}.png")
            with open(path, "wb") as f: f.write(r.content)
            local_charts[name] = path
            print(f"[DEBUG] Saved chart {name} to {path}", flush=True)
        except Exception as e:
            print(f"[ERROR] Failed to download chart '{name}': {e}", flush=True)

    # Build placeholder mapping
    placeholders = {}
    # Core URL & metadata fields
    for field in ["session_id", "report_date", "table_of_contents", "hw_gap_url", "sw_gap_url"]:
        if field in data:
            placeholders[f"{{{{ {field} }}}}"] = str(data[field])
    # Section titles
    for i in range(1, 21):
        placeholders[f"{{{{ section_{i}_title }}}}"] = data.get(f"section_{i}_title", "")
    # Content blocks
    for i in range(1, 21):
        placeholders[f"{{{{ content_{i} }}}}"] = data.get(f"content_{i}", "")

    # Slide narratives
    slide_keys = [
        'executive_summary', 'it_landscape_overview',
        'hardware_analysis', 'software_analysis', 'tier_classification_summary',
        'hardware_lifecycle_chart', 'software_licensing_review', 'security_vulnerability_heatmap',
        'performance_&_uptime_trends', 'system_reliability_overview', 'scalability_insights',
        'legacy_system_exposure', 'obsolete_platform_matrix', 'cloud_migration_targets',
        'strategic_it_alignment', 'business_impact_of_gaps', 'cost_of_obsolescence',
        'environmental_impact_and_sustainability',        # fixed key to match mapping :contentReference[oaicite:3]{index=3}
        'remediation_recommendations', 'roadmap_&_next_steps'
    ]
    for key in slide_keys:
        placeholders[f"{{{{ slide_{key} }}}}"] = str(data.get(f"slide_{key}", ""))

    # --------- DOCX Generation ---------
    doc = Document(TEMPLATE_DOCX)
    for para in doc.paragraphs:
        for ph, val in placeholders.items():
            if ph in para.text:
                para.text = para.text.replace(ph, val)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for ph, val in placeholders.items():
                    if ph in cell.text:
                        cell.text = cell.text.replace(ph, val)
    # Append charts at end of DOCX
    for chart_path in local_charts.values():
        doc.add_page_break()
        doc.add_picture(chart_path, width=Inches(6))
    docx_filename = f"IT_Current_Status_Assessment_Report_{session_id}.docx"
    docx_out = os.path.join(session_dir, docx_filename)
    doc.save(docx_out)
    print(f"[DEBUG] Saved DOCX to: {docx_out}", flush=True)

    # --------- PPTX Generation ---------
    prs = Presentation(TEMPLATE_PPTX)
    # Replace text placeholders
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for ph, val in placeholders.items():
                    if ph in shape.text:
                        shape.text = shape.text.replace(ph, val)
    # Inject chart images onto their corresponding slides :contentReference[oaicite:4]{index=4}
    for name, chart_path in local_charts.items():
        if name in section_to_slide_map:
            idx = section_to_slide_map[name] - 1
            if 0 <= idx < len(prs.slides):
                slide = prs.slides[idx]
                slide.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(6))
    pptx_filename = f"IT_Current_Status_Executive_Report_{session_id}.pptx"
    pptx_out = os.path.join(session_dir, pptx_filename)
    prs.save(pptx_out)
    print(f"[DEBUG] Saved PPTX to: {pptx_out}", flush=True)

    # --------- Upload to Google Drive ---------
    docx_url = upload_to_drive(docx_out, docx_filename, session_id)
    print(f"[DEBUG] Uploaded DOCX to Drive: {docx_url}", flush=True)
    try:
        pptx_url = upload_to_drive(pptx_out, pptx_filename, session_id)
        print(f"[DEBUG] Uploaded PPTX to Drive: {pptx_url}", flush=True)
    except Exception as e:
        print(f"[ERROR] PPTX upload failed: {e}", flush=True)
        pptx_url = ""

    return {"docx_url": docx_url, "pptx_url": pptx_url}
